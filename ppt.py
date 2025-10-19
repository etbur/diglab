from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

def add_slide(title, content):
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Title Slide
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Humaniot Robotic Coder"
slide.placeholders[1].text = (
    "An Autonomous System for Full-Cycle Software Development and Deployment\n"
    "National AI Innovation Challenge\n"
    "August 2025"
)

# Problem Slide (Updated)
add_slide("The Problem", 
    "- Software development is time-consuming, requiring coordination across developers, testers, and DevOps teams.\n"
    "- In emerging markets, businesses face limited access to skilled developers and infrastructure.\n"
    "- Manual coding and deployment introduce errors, delays, and scalability challenges.\n"
    "- Security and quality control are inconsistent, especially in freelance-driven work.\n"
    "- Robotics and automation remain underutilized in software engineering."
)

# AI Solution
add_slide("The AI Solution – Humaniot Robotic Coder",
    "- Uses NLP to interpret human instructions.\n"
    "- Autonomously generates web, desktop, and mobile apps.\n"
    "- Deploys to cloud, servers, or app stores without manual input.\n"
    "- Robotic interface interacts with tools and environments.\n"
    "- Continuously tests, monitors, and updates deployed applications."
)

# Market Potential
add_slide("Market Potential",
    "**Ethiopia:**\n- Growing tech sector with strong demand for software.\n- AI strategy and innovation support by government.\n\n"
    "**Africa:**\n- $5B+ software market; AI to contribute $1.2T to GDP by 2030.\n\n"
    "**Global:**\n- Global software market expected to reach $1T by 2030.\n- DevOps tools: $20B+ market by 2027."
)

# Team Composition
add_slide("Team Composition",
    "- Lead AI Engineer – Expert in NLP and AI planning.\n"
    "- Software Architect – Full-stack development and system design.\n"
    "- Robotics Specialist – Physical and virtual interface automation.\n"
    "- DevOps Engineer – Cloud deployment and infrastructure.\n"
    "- Business Lead – Strategy, partnerships, and outreach."
)

# AI Startup Center Alignment
add_slide("Alignment with AI Startup Center",
    "- Aligned with Ethiopia’s national AI strategy.\n"
    "- Supports AI Startup Center’s mission of innovation and job creation.\n"
    "- Ready for incubation, mentorship, and pilot programs.\n"
    "- Scalable across sectors like education, health, and fintech."
)

# Conclusion
add_slide("Conclusion",
    "Humaniot Robotic Coder solves inefficiencies in software development:\n"
    "- Eliminates delays and improves quality through autonomous AI.\n"
    "- Delivers secure, scalable solutions without human bottlenecks.\n"
    "- It’s more than a tool — it’s the future of intelligent software delivery."
)

prs.save("Humaniot_Robotic_Coder_AI_Competition_Updated_Pitch1.pptx")
